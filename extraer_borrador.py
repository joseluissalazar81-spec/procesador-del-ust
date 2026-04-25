"""
extraer_borrador.py — Extrae texto de archivos borrador del docente
Soporta: .docx, .pptx, .xlsx, .pdf, .txt
Procesador DEL · UST 2026-1
"""
from __future__ import annotations
from io import BytesIO


def extraer_texto(file_bytes: bytes, filename: str) -> str:
    """
    Extrae texto plano de un archivo borrador del docente.

    Parámetros
    ----------
    file_bytes : contenido del archivo
    filename   : nombre del archivo (para detectar extensión)

    Retorna
    -------
    Texto extraído (máx ~8000 caracteres para no saturar el contexto LLM)
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    try:
        if ext == "docx":
            return _desde_docx(file_bytes)
        elif ext == "pptx":
            return _desde_pptx(file_bytes)
        elif ext == "xlsx":
            return _desde_xlsx(file_bytes)
        elif ext == "pdf":
            return _desde_pdf(file_bytes)
        elif ext in ("txt", "md"):
            return file_bytes.decode("utf-8", errors="ignore")[:8000]
        else:
            return f"[Formato .{ext} no soportado para extracción automática]"
    except Exception as e:
        return f"[Error al leer el archivo: {str(e)[:120]}]"


def _desde_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(BytesIO(data))
    partes = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            partes.append(t)
    # Tablas
    for table in doc.tables:
        for row in table.rows:
            fila = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if fila:
                partes.append(fila)
    return "\n".join(partes)[:8000]


def _desde_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(BytesIO(data))
    partes = []
    for i, slide in enumerate(prs.slides, 1):
        textos_slide = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        textos_slide.append(t)
        if textos_slide:
            partes.append(f"[Diapositiva {i}]")
            partes.extend(textos_slide)
    return "\n".join(partes)[:8000]


def _desde_xlsx(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(BytesIO(data), data_only=True)
    partes = []
    for sname in wb.sheetnames:
        ws = wb[sname]
        partes.append(f"[Hoja: {sname}]")
        for row in ws.iter_rows(values_only=True):
            vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if vals:
                partes.append(" | ".join(vals))
    return "\n".join(partes)[:8000]


def _desde_pdf(data: bytes) -> str:
    import pdfplumber
    partes = []
    with pdfplumber.open(BytesIO(data)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                partes.append(t.strip())
    return "\n".join(partes)[:8000]
